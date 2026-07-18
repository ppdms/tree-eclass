from ..models import ExtractedDocument, ExtractedUnit, SourceMetadata
from .base import ExtractionError, ExtractionLimits, enforce_limits


def extract(path: str, source: SourceMetadata, limits: ExtractionLimits) -> ExtractedDocument:
    try:
        from pptx import Presentation
        deck = Presentation(path)
    except Exception as exc:
        raise ExtractionError(f"could not open PPTX: {exc}") from exc
    units = []
    for number, slide in enumerate(deck.slides, 1):
        parts, heading = [], None
        if slide.shapes.title and slide.shapes.title.text.strip():
            heading = slide.shapes.title.text.strip()
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                parts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                parts.extend(" | ".join(cell.text.strip() for cell in row.cells) for row in shape.table.rows)
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append("Speaker notes:\n" + notes)
        except (AttributeError, ValueError):
            pass
        units.append(ExtractedUnit("slide", str(number), "\n".join(parts), heading=heading))
    return enforce_limits(ExtractedDocument(source.display_name, "presentation", units,
                                             {"slide_count": len(deck.slides)}), limits)
