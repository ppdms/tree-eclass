import json
import os
import tempfile
import unittest
import zipfile
from unittest.mock import patch

from app.knowledge.extractors.archive import extract as extract_archive
from app.knowledge.extractors.base import ExtractionLimits, detect_source
from app.knowledge.extractors.docx import extract as extract_docx
from app.knowledge.extractors.html import extract as extract_html
from app.knowledge.extractors.notebook import extract as extract_notebook
from app.knowledge.extractors.pdf import extract as extract_pdf
from app.knowledge.extractors.pptx import extract as extract_pptx
from app.knowledge.extractors.text import extract as extract_text
from app.knowledge.extractors.xlsx import extract as extract_xlsx
from app.knowledge.models import SourceMetadata


class KnowledgeExtractorTests(unittest.TestCase):
    def setUp(self):
        self.temp = tempfile.TemporaryDirectory()
        self.source = SourceMetadata(1, "Course", "C", "/Course/item", None, "item", "hash")
        self.limits = ExtractionLimits(max_characters=100_000)

    def tearDown(self):
        self.temp.cleanup()

    def test_text_line_locators(self):
        path = os.path.join(self.temp.name, "greek.txt")
        with open(path, "w", encoding="utf-8") as handle:
            handle.write("πρώτη\nδεύτερη\n")
        result = extract_text(path, self.source, self.limits)
        self.assertEqual(result.units[0].locator_type, "line")
        self.assertEqual(result.units[0].locator_start, "1")

    def test_notebook_cell_locators(self):
        path = os.path.join(self.temp.name, "lesson.ipynb")
        with open(path, "w", encoding="utf-8") as handle:
            json.dump({"cells": [{"cell_type": "markdown", "source": ["# Topic"]},
                                  {"cell_type": "code", "source": ["x = 1"]}]}, handle)
        result = extract_notebook(path, self.source, self.limits)
        self.assertEqual([unit.locator_start for unit in result.units], ["1", "2"])

    def test_archive_rejects_traversal_and_indexes_text(self):
        path = os.path.join(self.temp.name, "items.zip")
        with zipfile.ZipFile(path, "w") as archive:
            archive.writestr("safe/code.py", "print('course evidence')")
            archive.writestr("../escape.txt", "must not appear")
        result = extract_archive(path, self.source, self.limits)
        self.assertEqual([unit.locator_start for unit in result.units], ["safe/code.py"])
        self.assertTrue(any("unsafe" in warning for warning in result.warnings))

    def test_archive_indexes_notebook_members(self):
        path = os.path.join(self.temp.name, "notebooks.zip")
        notebook = {"cells": [{"cell_type": "markdown", "source": ["# Lab evidence"]},
                               {"cell_type": "code", "source": ["answer = 42"]}]}
        with zipfile.ZipFile(path, "w") as archive:
            archive.writestr("labs/lab1.ipynb", json.dumps(notebook))
            archive.writestr("__MACOSX/._lab1.ipynb", b"AppleDouble metadata, not a notebook")
        result = extract_archive(path, self.source, self.limits)
        self.assertEqual(
            [unit.locator_start for unit in result.units],
            ["labs/lab1.ipynb#1", "labs/lab1.ipynb#2"],
        )
        self.assertIn("Lab evidence", result.units[0].text)

    def test_pdf_page_locators_are_one_based(self):
        from pypdf import PdfWriter
        path = os.path.join(self.temp.name, "slides.pdf")
        writer = PdfWriter()
        writer.add_blank_page(width=300, height=300)
        writer.add_blank_page(width=300, height=300)
        with open(path, "wb") as handle:
            writer.write(handle)
        result = extract_pdf(path, self.source, self.limits)
        self.assertEqual([unit.locator_start for unit in result.units], ["1", "2"])

    def test_content_detection_prefers_pdf_magic_bytes_and_rejects_html_downloads(self):
        kind, mime, reason = detect_source("opaque-download", "text/html", b"%PDF-1.7\nbody")
        self.assertEqual((kind, mime, reason), ("pdf", "application/pdf", None))

        kind, mime, reason = detect_source("slides.pdf", "application/pdf", b"<!doctype html><html>")
        self.assertEqual((kind, mime, reason), (None, "text/html", "download_html"))

        kind, mime, reason = detect_source("/Course/regression.slides.html", "text/plain",
                                           b"<!doctype html><html><p>Evidence</p>")
        self.assertEqual((kind, mime, reason), ("html", "text/html", None))

    def test_sparse_pdf_page_uses_ocr_with_provenance(self):
        from pypdf import PdfWriter
        path = os.path.join(self.temp.name, "scan.pdf")
        writer = PdfWriter()
        writer.add_blank_page(width=300, height=300)
        with open(path, "wb") as handle:
            writer.write(handle)
        limits = ExtractionLimits(max_characters=100_000, ocr_enabled=True, ocr_languages="ell+eng")
        with patch("app.knowledge.extractors.pdf._poppler_page", return_value=""), \
             patch("app.knowledge.extractors.pdf._ocr_page", return_value=(
                 "OCR extracted course evidence", {"engine": "tesseract", "languages": "ell+eng", "dpi": "200"}
             )) as ocr:
            result = extract_pdf(path, self.source, limits)
        ocr.assert_called_once()
        self.assertEqual(result.units[0].text, "OCR extracted course evidence")
        self.assertEqual(result.units[0].metadata["provenance"], "tesseract_ocr")
        self.assertEqual(result.units[0].metadata["ocr"]["engine"], "tesseract")

    def test_office_and_html_locators(self):
        from docx import Document
        from openpyxl import Workbook
        from pptx import Presentation

        docx_path = os.path.join(self.temp.name, "notes.docx")
        document = Document()
        document.add_heading("Definition", level=1)
        document.add_paragraph("A bounded definition.")
        document.save(docx_path)
        docx_result = extract_docx(docx_path, self.source, self.limits)
        self.assertEqual(docx_result.units[0].locator_type, "section")
        self.assertEqual(docx_result.units[0].heading, "Definition")

        pptx_path = os.path.join(self.temp.name, "deck.pptx")
        deck = Presentation()
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = "First slide"
        slide.placeholders[1].text = "Evidence"
        deck.save(pptx_path)
        pptx_result = extract_pptx(pptx_path, self.source, self.limits)
        self.assertEqual(pptx_result.units[0].locator_start, "1")

        xlsx_path = os.path.join(self.temp.name, "data.xlsx")
        workbook = Workbook()
        workbook.active.title = "Exercises"
        workbook.active["A1"] = "Question"
        workbook.save(xlsx_path)
        xlsx_result = extract_xlsx(xlsx_path, self.source, self.limits)
        self.assertEqual(xlsx_result.units[0].locator_start, "Exercises")

        html_path = os.path.join(self.temp.name, "page.html")
        with open(html_path, "w", encoding="utf-8") as handle:
            handle.write("<html><title>Lesson</title><h1>Topic</h1><p>Evidence</p><script>ignore()</script></html>")
        html_result = extract_html(html_path, self.source, self.limits)
        self.assertEqual(html_result.title, "Lesson")
        self.assertEqual(html_result.units[0].heading, "Topic")


if __name__ == "__main__":
    unittest.main()
